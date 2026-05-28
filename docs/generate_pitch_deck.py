"""Generate Agent Checkout pitch deck (PowerPoint)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ACCENT = RGBColor(0, 120, 80)
DARK = RGBColor(30, 30, 30)
GRAY = RGBColor(90, 90, 90)


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.LEFT
    sub = tf.add_paragraph()
    sub.text = subtitle
    sub.font.size = Pt(18)
    sub.font.color.rgb = GRAY
    sub.space_before = Pt(12)


def add_section_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], note: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.level = 0
        p.space_after = Pt(10)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.8), Inches(0.8))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(11)
        np.font.italic = True
        np.font.color.rgb = GRAY


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left: list[str], right_title: str, right: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True

    for col, col_title, items, x in [(0, left_title, left, 0.6), (1, right_title, right, 5.2)]:
        h = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.2), Inches(0.5))
        h.text_frame.paragraphs[0].text = col_title
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.size = Pt(18)
        h.text_frame.paragraphs[0].font.color.rgb = ACCENT
        b = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(4.2), Inches(5))
        tf = b.text_frame
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.space_after = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Agent Checkout",
        "Agent-native Shopify commerce on Somnia\nEncode Agentathon 2026",
    )

    add_bullet_slide(
        prs,
        "Executive Summary",
        [
            "Agent Checkout makes any Shopify store purchasable by AI agents in ~2 minutes",
            "Agents discover products via MCP, pay via HTTP 402 + native STT on Somnia",
            "Somnia Agents power pricing oracle, address validation, and NL product search",
            "Confirmed orders sync to Shopify Admin — merchants keep existing workflows",
            "Open-source prototype deployed on Shannon Testnet (chain ID 50312)",
        ],
    )

    add_section_slide(prs, "The Problem")

    add_bullet_slide(
        prs,
        "Commerce Was Built for Humans, Not Agents",
        [
            "4.8M+ Shopify merchants — none have a standard agent checkout interface",
            "AI agents cannot reliably browse, cart, and pay across arbitrary storefronts",
            "Payment stacks (Stripe, cards) assume human authorization and sessions",
            "Custom agent integrations are expensive, fragile, and non-composable",
            "Merchants miss a new distribution channel: autonomous AI shoppers",
        ],
        "Sources: Shopify merchant stats; agent payment protocol adoption reports 2025–2026",
    )

    add_section_slide(prs, "Market Opportunity")

    add_bullet_slide(
        prs,
        "Agentic Commerce Is Becoming Infrastructure",
        [
            "WEF projection: ~$30T agent economy potential by 2030",
            "AI agent market: ~$7.8B (2026) → ~$52B projected (2030)",
            "12,000+ production MCP servers by mid-2026 (up from ~100 in early 2025)",
            "x402: 50M+ transactions; Linux Foundation standard (Coinbase, Cloudflare, Shopify participants)",
            "Visa TAP + Mastercard Agent Pay: 5,000+ merchants in pilot programs",
            "18% of B2B SaaS already deploy MCP servers; ecommerce catching up fast",
        ],
        "Sources: Presenc AI Agentic Commerce Benchmarks 2026; primores.org agent payment wiki",
    )

    add_two_column_slide(
        prs,
        "Market Gap We Fill",
        "Existing solutions",
        [
            "x402 rails — payment only, no merchant onboarding",
            "Shopify UCP/AP2 — authorization schema, not agent checkout",
            "Generic MCP servers — no payments or order fulfillment",
            "Solana x402 demos — not Shopify-native, not Somnia",
        ],
        "Agent Checkout",
        [
            "Full loop: discover → 402 → pay → verify → Shopify order",
            "2-minute merchant setup, no code changes",
            "Dual MCP architecture (shop + pay)",
            "Somnia-native settlement + Somnia Agents",
        ],
    )

    add_section_slide(prs, "Our Solution")

    add_bullet_slide(
        prs,
        "Product: Agent Checkout",
        [
            "Merchant UI (Next.js): register store, select products, monitor orders",
            "Shopping MCP (:3001): list_stores, search_products_nl, initiate/finalize checkout",
            "Payment MCP (:3002): make_stt_payment with on-chain proof",
            "HTTP 402 two-phase checkout with x402-somnia-stt-v1 scheme",
            "Supabase for stores, products, order intents, orders",
            "Shopify Admin API creates paid orders after on-chain verification",
        ],
    )

    add_bullet_slide(
        prs,
        "Checkout Flow (HTTP 402)",
        [
            "Phase 1: Agent calls initiate_checkout → server returns HTTP 402 + STT payment requirements",
            "Somnia Agents validate address + fetch STT/USD rate before quoting",
            "Agent pays native STT via Payment MCP → receives txHash proof",
            "Phase 2: finalize_checkout with X-PAYMENT header → viem verifies receipt",
            "Server creates paid order in Shopify → merchant sees it in dashboard + admin",
        ],
    )

    add_section_slide(prs, "Why Somnia")

    add_bullet_slide(
        prs,
        "Somnia: The Agentic L1",
        [
            "Purpose-built L1 for agents and real-time applications — 1M+ TPS, sub-second finality",
            "Somnia Agents: consensus-validated off-chain compute (JSON API, LLM inference)",
            "Native STT/SOMI settlement — aligns with micropayment economics (sub-cent fees)",
            "EVM-compatible — viem, Hardhat, existing Solidity tooling",
            "Reactive architecture — smart contracts can respond to on-chain events in real time",
            "Encode Agentathon + Somnia AI Hackathon alignment — ecosystem actively funding agent apps",
        ],
        "Sources: somnia.network, docs.somnia.network/agents",
    )

    add_two_column_slide(
        prs,
        "Why Somnia vs. Other Chains",
        "Alternatives",
        [
            "Base/Solana x402 — mature rails but no native agent compute layer",
            "Ethereum L1 — high fees break agent micropayments",
            "Centralized APIs — no verifiable consensus on oracle/LLM outputs",
        ],
        "Somnia advantage",
        [
            "Agents + payments + commerce on one chain",
            "Deterministic LLMs enable consensus on AI outputs",
            "Hyper-performance for high-frequency agent transactions",
            "Strong hackathon/ecosystem signal for hiring & grants",
        ],
    )

    add_section_slide(prs, "Unique Selling Points")

    add_bullet_slide(
        prs,
        "What Makes Agent Checkout Different",
        [
            "Only Shopify + MCP + HTTP 402 + Somnia Agents in one open-source stack",
            "Merchants: zero blockchain knowledge — 2-minute setup",
            "Agents: standard MCP tools — no scraping or custom per-store APIs",
            "Payments: verifiable on-chain STT, not trust-me callbacks",
            "Intelligence: Somnia Agents for oracle, validation, NL search (not just chain swap)",
            "Orders: land in Shopify Admin like any normal sale",
        ],
    )

    add_section_slide(prs, "What We Have Built")

    add_bullet_slide(
        prs,
        "Development Progress (Complete Prototype)",
        [
            "Migrated from Solana/USDC template → Somnia STT + viem SDK",
            "Cursor dev environment: rules, skills, subagents for Somnia integration",
            "packages/backend: Express API, dual MCP servers, x402-sdk, somnia-agents module",
            "packages/frontend: Agent Checkout branded UI, dashboard with explorer links",
            "Security: body-hash anti-tampering, 15-min intent expiry, receipt verification",
            "Docs: comprehensive README, SUBMISSION.md, this pitch deck",
        ],
    )

    add_bullet_slide(
        prs,
        "Technology Stack",
        [
            "Frontend: Next.js 16, React 19, Tailwind CSS 4",
            "Backend: Express, TypeScript, Zod, Model Context Protocol SDK",
            "Database: Supabase (PostgreSQL)",
            "Blockchain: Somnia Shannon Testnet, viem, native STT",
            "Commerce: Shopify Admin REST API",
            "Protocols: HTTP 402 (x402), MCP JSON-RPC 2.0",
        ],
    )

    add_section_slide(prs, "Business Model & GTM")

    add_bullet_slide(
        prs,
        "Monetization Paths",
        [
            "SaaS: monthly fee per agent-enabled Shopify store",
            "Take rate: small % on agent-settled STT volume",
            "Enterprise: white-label MCP + compliance for large merchants",
            "Protocol services: hosted facilitators, oracle, and agent invocation",
            "Ecosystem grants: Somnia, Encode, Shopify partner programs",
        ],
    )

    add_bullet_slide(
        prs,
        "Go-to-Market",
        [
            "Phase 1: Hackathon demo + open-source GitHub (now)",
            "Phase 2: 10 pilot Shopify merchants on Shannon testnet",
            "Phase 3: Somnia mainnet + Shopify App Store listing",
            "Phase 4: Agent marketplace — agents discover all Agent Checkout stores",
            "Community: Discord dev-rel, demo video, blog post on Somnia Agents + x402",
        ],
    )

    add_section_slide(prs, "Competition & Moat")

    add_two_column_slide(
        prs,
        "Competitive Landscape",
        "Direct / adjacent",
        [
            "Coinbase x402 demos — payment rail only",
            "Google AP2 + UCP — authorization, early merchant coalition",
            "Stripe MPP — HTTP auth scheme, Stripe ecosystem",
            "Custom agent wrappers per merchant — fragile",
        ],
        "Our moat",
        [
            "Shopify-native merchant UX + agent MCP in one product",
            "Somnia Agents integration depth (not generic EVM port)",
            "Open-source reference implementation for agent commerce",
            "First-mover in 'Shopify agent checkout on Somnia'",
        ],
    )

    add_section_slide(prs, "Roadmap")

    add_bullet_slide(
        prs,
        "Product Roadmap",
        [
            "Q2 2026: Encode Agentathon submission, demo video, deploy to Vercel/Railway",
            "Q3 2026: Production Somnia Agents (JSON API + LLM Inference on-chain)",
            "Q3 2026: Somnia Reactivity — auto-fulfill on payment events",
            "Q4 2026: Mainnet SOMI, Shopify App, multi-store agent catalog",
            "2027: AP2/UCP mandate compatibility, enterprise compliance layer",
        ],
    )

    add_section_slide(prs, "Team & Ask")

    add_bullet_slide(
        prs,
        "Why We Will Win",
        [
            "Shipping full-stack prototype, not slides-only",
            "Deep alignment with Somnia's agentic L1 thesis",
            "Timing: agent payment protocols hitting production adoption in 2026",
            "Shopify TAM: millions of merchants need agent readiness",
            "Clear demo narrative for judges and investors",
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "Agent Checkout — Where AI agents shop on Somnia\nGitHub: agent-checkout  |  Encode Agentathon 2026",
    )

    out = "d:/agent-checkout/docs/Agent_Checkout_Pitch_Deck.pptx"
    prs.save(out)
    print(f"Created: {out}")


if __name__ == "__main__":
    main()
